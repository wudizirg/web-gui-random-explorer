"""
生成「软件测试导论 期末大作业」完善方案 PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 颜色方案 ──
COLOR_PRIMARY = RGBColor(0x1A, 0x56, 0xDB)   # 深蓝
COLOR_SECONDARY = RGBColor(0x10, 0xB9, 0x81)  # 绿色
COLOR_ACCENT = RGBColor(0xF5, 0x9E, 0x0B)     # 橙色
COLOR_DARK = RGBColor(0x1E, 0x29, 0x3B)        # 深灰黑
COLOR_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)       # 浅白
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY = RGBColor(0x64, 0x74, 0x8B)
COLOR_RED = RGBColor(0xEF, 0x44, 0x44)
COLOR_CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# ── 辅助函数 ──
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=COLOR_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=COLOR_DARK, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = spacing
    return txBox

def add_code_block(slide, left, top, width, height, code_text, font_size=12):
    shape = add_rect(slide, left, top, width, height, COLOR_CODE_BG)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0xA6, 0xE2, 0x2E)  # 终端绿
    p.font.name = "Consolas"
    p.alignment = PP_ALIGN.LEFT
    return shape

def add_step_card(slide, left, top, width, height, step_num, title, desc, color):
    """添加步骤卡片"""
    # 卡片背景
    card = add_rect(slide, left, top, width, height, COLOR_WHITE)
    card.shadow.inherit = False
    # 左侧色条
    add_rect(slide, left, top, Inches(0.08), height, color)
    # 步骤编号圆圈
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.25), top + Inches(0.2), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(step_num)
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # 标题
    add_text_box(slide, left + Inches(0.85), top + Inches(0.18), width - Inches(1.1), Inches(0.4),
                 title, font_size=20, color=color, bold=True)
    # 描述
    add_text_box(slide, left + Inches(0.3), top + Inches(0.75), width - Inches(0.6), height - Inches(1.0),
                 desc, font_size=13, color=COLOR_GRAY)

# ═══════════════════════════════════════════
# Slide 1: 封面
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, COLOR_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), COLOR_PRIMARY)
add_rect(slide, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), COLOR_PRIMARY)

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.0),
             "Web GUI 自动化探索测试工具", font_size=44, color=COLOR_WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.8),
             "Abtext — 软件测试导论 期末大作业 完善方案", font_size=28, color=COLOR_SECONDARY, bold=False)
add_rect(slide, Inches(1.5), Inches(3.7), Inches(3), Inches(0.04), COLOR_SECONDARY)

items = [
    "方向一：增加页面状态抽象与表单结构感知",
    "方向二：接入大模型，生成更聪明的点击/输入序列",
    "方向三：增加截图、DOM 快照、异常检测与回溯",
]
add_bullet_list(slide, Inches(1.5), Inches(4.1), Inches(10), Inches(2.5), items,
                font_size=20, color=RGBColor(0x94, 0xA3, 0xB8))

add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5),
             "2026年6月  |  TypeScript + Playwright + LLM", font_size=16, color=COLOR_GRAY)

# ═══════════════════════════════════════════
# Slide 2: 目录
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_LIGHT)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), COLOR_PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "📋 目  录", font_size=36, color=COLOR_WHITE, bold=True)

toc_items = [
    ("01", "项目背景与现状分析", "了解 Abtext 的当前架构与能力边界"),
    ("02", "方向一：页面状态抽象与表单感知", "FSM 状态机 + 表单结构解析"),
    ("03", "方向二：大模型智能策略增强", "LLM 多步规划 + 状态感知 Prompt"),
    ("04", "方向三：截图/DOM快照/异常检测/回溯", "全链路监控与故障恢复"),
    ("05", "整体架构设计", "新模块关系图与数据流"),
    ("06", "实施计划与产出", "10步实施路线 + 预期成果"),
    ("07", "修改前后详细对比", "文件结构、代码量、能力维度"),
    ("08", "实际完成情况", "三个方向的具体实现成果"),
    ("09", "后续完善与展望", "短/中/长期优化路线图"),
]
for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.6) + Inches(i * 0.9)
    add_text_box(slide, Inches(1.2), y, Inches(0.6), Inches(0.5),
                 num, font_size=28, color=COLOR_PRIMARY, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.0), Inches(8), Inches(0.45),
                 title, font_size=22, color=COLOR_DARK, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.45), Inches(8), Inches(0.35),
                 desc, font_size=14, color=COLOR_GRAY)

# ═══════════════════════════════════════════
# Slide 3: 项目背景
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_LIGHT)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), COLOR_PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "01  项目背景与现状分析", font_size=32, color=COLOR_WHITE, bold=True)

# 左侧：项目简介
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
             "📦 Abtext 是什么？", font_size=22, color=COLOR_PRIMARY, bold=True)
intro = [
    "面向网页 GUI 测试输入生成的大作业项目",
    "基于 Playwright 实现浏览器自动化控制",
    "支持随机探索 + LLM（OpenAI）辅助决策",
    "每步自动截图、记录动作日志",
    "TypeScript 编写，Node.js 运行",
]
add_bullet_list(slide, Inches(0.8), Inches(1.95), Inches(5.5), Inches(2.5), intro, font_size=15, color=COLOR_DARK)

# 右侧：当前架构
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5.5), Inches(0.5),
             "🏗️ 当前架构", font_size=22, color=COLOR_PRIMARY, bold=True)
code = (
    "Abtext-main/\n"
    "├── src/\n"
    "│   ├── cli.ts        ← CLI 入口 & 参数解析\n"
    "│   ├── explorer.ts   ← 核心引擎\n"
    "│   │   • annotateAndCollectElements()\n"
    "│   │   • capturePageContext() 截图\n"
    "│   │   • callLlm() LLM 调用\n"
    "│   │   • executeAction() 动作执行\n"
    "│   │   • fallbackRandomStep() 随机回退\n"
    "│   └── types.ts      ← 类型定义\n"
    "├── screenshots/      ← 截图输出\n"
    "└── package.json"
)
add_code_block(slide, Inches(7.2), Inches(1.95), Inches(5.5), Inches(3.5), code)

# 底部：已具备 vs 缺失
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(0.5),
             "✅ 现有能力", font_size=18, color=COLOR_SECONDARY, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(5.25), Inches(5.5), Inches(2.0),
    ["Playwright 浏览器自动化", "随机点击/输入/滚动", "LLM 视觉辅助决策（截图输入）",
     "元素标注追踪", "多轮对话式交互"], font_size=14)

add_text_box(slide, Inches(7.2), Inches(4.8), Inches(5.5), Inches(0.5),
             "❌ 当前缺失", font_size=18, color=COLOR_RED, bold=True)
add_bullet_list(slide, Inches(7.2), Inches(5.25), Inches(5.5), Inches(2.0),
    ["无页面状态建模（FSM）", "无表单结构感知", "无 DOM 快照保存",
     "无异常检测机制", "无状态回溯能力"], font_size=14, color=COLOR_RED)

# ═══════════════════════════════════════════
# Slide 4: 方向一 概览
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_LIGHT)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), COLOR_SECONDARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "02  方向一：页面状态抽象与表单结构感知", font_size=32, color=COLOR_WHITE, bold=True)

# 左侧：状态机
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(6.0), Inches(0.5),
             "🔁 页面状态机 (FSM)", font_size=22, color=COLOR_SECONDARY, bold=True)

fsm_items = [
    "页面指纹 = URL 路径 + 关键元素签名（form action, input name 集合）",
    "自动检测状态转移：导航 → 新 URL → 计算指纹 → 判断已知/未知状态",
    "状态分类：InitialPage / LoginPage / FormPage / ContentPage / ErrorPage",
    "记录完整的状态转移图（graph），支持事后可视化分析",
]
add_bullet_list(slide, Inches(0.8), Inches(1.95), Inches(6.0), Inches(2.2), fsm_items, font_size=14)

# FSM 示意图用代码块
fsm_code = (
    "// state-machine.ts - 核心接口\n"
    "interface PageFingerprint {\n"
    "  urlPath: string;\n"
    "  formSignatures: string[];  // form action + input names\n"
    "  keyElements: string[];     // 关键元素 hash\n"
    "}\n"
    "class PageStateMachine {\n"
    "  fingerprint(page): PageFingerprint\n"
    "  classify(fingerprint): PageState\n"
    "  transition(from, to, action)\n"
    "  isNewState(fingerprint): boolean\n"
    "}"
)
add_code_block(slide, Inches(0.8), Inches(4.3), Inches(6.0), Inches(2.8), fsm_code, font_size=11)

# 右侧：表单分析
add_text_box(slide, Inches(7.5), Inches(1.4), Inches(5.5), Inches(0.5),
             "📝 表单结构感知", font_size=22, color=COLOR_SECONDARY, bold=True)

form_items = [
    "识别 <form> 标签边界与嵌套关系",
    "解析 <label for='...'> 关联，提取字段语义名",
    "检测 required / aria-required / pattern 约束",
    "识别输入类型：text / password / email / tel / number",
    "分组同表单字段，按 DOM 顺序排列",
    "输出结构化 FormSchema 供 LLM 使用",
]
add_bullet_list(slide, Inches(7.5), Inches(1.95), Inches(5.2), Inches(2.5), form_items, font_size=14)

form_code = (
    "// form-analyzer.ts - 核心接口\n"
    "interface FormField {\n"
    "  name: string;        // 语义名（label 或 name）\n"
    "  type: string;        // text/password/email...\n"
    "  required: boolean;\n"
    "  selector: string;\n"
    "  filled: boolean;\n"
    "}\n"
    "interface FormSchema {\n"
    "  fields: FormField[];\n"
    "  submitSelector?: string;\n"
    "}\n"
    "class FormAnalyzer {\n"
    "  analyze(page): FormSchema[]\n"
    "  getFormContext(page): string  // LLM prompt\n"
    "}"
)
add_code_block(slide, Inches(7.5), Inches(4.3), Inches(5.5), Inches(2.8), form_code, font_size=11)

# ═══════════════════════════════════════════
# Slide 5: 方向一 状态转移图
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_LIGHT)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), COLOR_SECONDARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "方向一  页面状态转移模型", font_size=32, color=COLOR_WHITE, bold=True)

# Mermaid 风格的 ASCII 状态图用代码块展示
state_diagram = (
    "                    ┌──────────────────────────────┐\n"
    "                    │        [*] 开始探索            │\n"
    "                    └─────────────┬────────────────┘\n"
    "                                  │\n"
    "                    ┌─────────────▼────────────────┐\n"
    "                    │      InitialPage             │\n"
    "                    │  (页面指纹 = URL + 元素签名)    │\n"
    "                    └──┬──────────┬────────────┬───┘\n"
    "                       │          │            │\n"
    "           检测到登录   │   检测到  │   无表单    │\n"
    "           表单        │   普通表单 │            │\n"
    "          ┌────────────▼──┐ ┌──────▼──────┐ ┌──▼────────┐\n"
    "          │  LoginPage    │ │  FormPage   │ │ContentPage│\n"
    "          └───┬───────┬───┘ └──┬──────┬───┘ └───────────┘\n"
    "              │       │        │      │\n"
    "         登录成功  登录失败  提交成功  验证失败\n"
    "              │       │        │      │\n"
    "     ┌────────▼──┐ ┌─▼──────────┐ ┌──▼──────────┐\n"
    "     │LoggedInPage│ │LoginError  │ │FormSubmitted │\n"
    "     └────────────┘ │   Page     │ │    Page      │\n"
    "                    └─────┬──────┘ └──────────────┘\n"
    "                          │ 自动重试\n"
    "                          ▼\n"
    "                    回到 LoginPage"
)
add_code_block(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(5.5), state_diagram, font_size=15)

# ═══════════════════════════════════════════
# Slide 6: 方向二 概览
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_LIGHT)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), COLOR_ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "03  方向二：接入大模型，生成更聪明的点击/输入序列", font_size=32, color=COLOR_WHITE, bold=True)

# 三列卡片
cards = [
    ("🧠", "表单结构注入 Prompt", [
        "将 FormAnalyzer 分析结果",
        "注入 LLM 系统提示词",
        "告知字段名、类型、必填状态",
        "\"此表单: 用户名(text,required)",
        "→密码(password,required)→登录\"",
    ]),
    ("📋", "多步规划 (Plan-then-Act)", [
        "LLM 先输出 2-5 步计划",
        "再逐步执行每个子动作",
        "执行过程中验证每步结果",
        "若步骤失败则重新规划",
        "避免「一步走错全盘皆输」",
    ]),
    ("🔄", "状态感知对话", [
        "每步将当前 FSM 状态告知 LLM",
        "已访问状态列表 → 避免重复",
        "当前表单填充进度",
        "\"第2步: LoginPage 状态",
        "已填:用户名, 待填:密码\"",
    ]),
]
for i, (icon, title, items) in enumerate(cards):
    left = Inches(0.6 + i * 4.15)
    add_rect(slide, left, Inches(1.5), Inches(3.85), Inches(4.0), COLOR_WHITE)
    add_text_box(slide, left + Inches(0.2), Inches(1.65), Inches(3.4), Inches(0.5),
                 f"{icon} {title}", font_size=18, color=COLOR_ACCENT, bold=True)
    add_bullet_list(slide, left + Inches(0.2), Inches(2.2), Inches(3.4), Inches(3.0),
                    items, font_size=13, spacing=Pt(6))

# 底部代码
add_text_box(slide, Inches(0.8), Inches(5.7), Inches(6.0), Inches(0.4),
             "新增文件: src/llm-strategist.ts", font_size=16, color=COLOR_PRIMARY, bold=True)
strategy_code = (
    "class LlmStrategist {\n"
    "  buildSmartPrompt(state, formSchemas, history): string\n"
    "  parseMultiStepPlan(raw: string): PlannedStep[]\n"
    "  validateStepExecution(planned, actual): boolean\n"
    "  shouldReplan(failedStep, context): boolean\n"
    "}"
)
add_code_block(slide, Inches(0.8), Inches(6.1), Inches(12), Inches(1.1), strategy_code, font_size=13)

# ═══════════════════════════════════════════
# Slide 7: 方向二 Prompt 增强对比
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_LIGHT)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), COLOR_ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "方向二  Prompt 增强对比", font_size=32, color=COLOR_WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
             "❌ 当前 Prompt（无结构感知）", font_size=18, color=COLOR_RED, bold=True)
old_prompt = (
    "Current page:\n"
    "- URL: https://example.com/login\n"
    "- Title: Login\n\n"
    "Numbered candidates:\n"
    "- [5] input — input type=text (input)\n"
    "- [6] input — input type=password (input)\n"
    "- [2] click — button \"提交\"\n\n"
    "Return ONE JSON..."
)
add_code_block(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(3.5), old_prompt, font_size=11)

add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5.5), Inches(0.5),
             "✅ 增强 Prompt（含结构感知）", font_size=18, color=COLOR_SECONDARY, bold=True)
new_prompt = (
    "Current page:\n"
    "- URL: https://example.com/login\n"
    "- Title: Login\n"
    "- State: LoginPage (首次访问)\n\n"
    "表单结构分析:\n"
    "┌─ Form#login-form ───────────┐\n"
    "│ 1. 用户名 [text, required]   │\n"
    "│ 2. 密码   [password,required]│\n"
    "│ 3. 记住我 [checkbox]         │\n"
    "│ [提交] 按钮                  │\n"
    "└─────────────────────────────┘\n\n"
    "填充进度: 0/2 必填字段完成\n"
    "建议: 先填 用户名 → 密码 → 点击提交\n\n"
    "Return ONE JSON..."
)
add_code_block(slide, Inches(7.2), Inches(1.9), Inches(5.5), Inches(3.5), new_prompt, font_size=11)

add_text_box(slide, Inches(0.8), Inches(5.7), Inches(12), Inches(0.4),
             "🎯 效果：LLM 能准确理解表单结构、按顺序填写、避免重复操作、一次性完成表单提交流程",
             font_size=16, color=COLOR_PRIMARY, bold=True)

# ═══════════════════════════════════════════
# Slide 8: 方向三 概览
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_LIGHT)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), RGBColor(0x8B, 0x5C, 0xF6))
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "04  方向三：截图、DOM 快照、异常检测与回溯", font_size=32, color=COLOR_WHITE, bold=True)

# 四个模块卡片
modules = [
    ("📸", "截图增强", "已有", [
        "每步自动截图 fullPage",
        "截图带步骤编号命名",
        "新增：异常时刻自动截图",
        "新增：截图与 DOM 快照关联索引",
    ]),
    ("🌳", "DOM 快照", "新增", [
        "每步保存简化 DOM 树(JSON)",
        "记录关键元素属性与状态",
        "支持事后离线回放分析",
        "增量对比：前后两步 DOM diff",
    ]),
    ("🚨", "异常检测", "新增", [
        "page.on('pageerror') JS 异常",
        "page.on('console') error 级别",
        "page.on('crash') 页面崩溃",
        "HTTP 4xx/5xx 状态关键词检测",
        "页面空白/白屏检测",
    ]),
    ("⏪", "状态回溯", "新增", [
        "保存每步页面状态快照",
        "异常触发自动 goBack()",
        "恢复到已知良好状态继续",
        "可配置最大回溯次数(默认3)",
        "回溯失败 → 优雅降级",
    ]),
]
for i, (icon, title, tag, items) in enumerate(modules):
    left = Inches(0.4 + i * 3.2)
    tag_color = COLOR_SECONDARY if tag == "已有" else COLOR_ACCENT
    add_rect(slide, left, Inches(1.5), Inches(3.0), Inches(4.5), COLOR_WHITE)
    add_text_box(slide, left + Inches(0.15), Inches(1.6), Inches(2.7), Inches(0.4),
                 f"{icon} {title}", font_size=18, color=RGBColor(0x8B, 0x5C, 0xF6), bold=True)
    # tag
    tag_shape = add_rect(slide, left + Inches(2.0), Inches(1.6), Inches(0.75), Inches(0.3), tag_color)
    tf = tag_shape.text_frame
    tf.paragraphs[0].text = tag
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_bullet_list(slide, left + Inches(0.15), Inches(2.1), Inches(2.7), Inches(3.5),
                    items, font_size=12, spacing=Pt(5))

# 底部
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(6.0), Inches(0.4),
             "新增文件: dom-snapshot.ts / anomaly-detector.ts / backtracker.ts", font_size=16, color=COLOR_PRIMARY, bold=True)

# ═══════════════════════════════════════════
# Slide 9: 方向三 异常检测流程
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_LIGHT)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), RGBColor(0x8B, 0x5C, 0xF6))
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "方向三  异常检测与回溯流程", font_size=32, color=COLOR_WHITE, bold=True)

flow = (
    "  每一步执行后:\n"
    "  ┌──────────────────────────────────────────────────────┐\n"
    "  │  1. 截图保存  ──→  2. DOM 快照  ──→  3. 异常检测      │\n"
    "  └──────────────────────────────────────────────────────┘\n"
    "                              │\n"
    "                     ┌───────▼────────┐\n"
    "                     │  检测到异常？    │\n"
    "                     └───┬────────┬───┘\n"
    "                         │ YES    │ NO\n"
    "                    ┌────▼───┐    │\n"
    "                    │记录异常  │    │\n"
    "                    │类型+截图 │    │\n"
    "                    └───┬────┘    │\n"
    "                   ┌────▼─────┐   │\n"
    "                   │ 回溯次数  │   │\n"
    "                   │ < 最大阈值│   │\n"
    "                   └──┬───┬──┘   │\n"
    "                  YES │   │ NO   │\n"
    "              ┌───────▼┐  │      │\n"
    "              │goBack()│  │      │\n"
    "              │恢复状态 │  │      │\n"
    "              └───┬───┘  │      │\n"
    "                  │      │      │\n"
    "            继续探索   跳过该步  继续探索\n"
    "                  │      │      │\n"
    "                  └──────┴──────┘\n"
    "                         │\n"
    "                    ┌────▼─────┐\n"
    "                    │ 下一轮循环│\n"
    "                    └──────────┘"
)
add_code_block(slide, Inches(2.0), Inches(1.4), Inches(9), Inches(5.8), flow, font_size=14)

# ═══════════════════════════════════════════
# Slide 10: 整体架构设计
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_LIGHT)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), COLOR_PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "05  整体架构设计", font_size=32, color=COLOR_WHITE, bold=True)

arch = (
    "                          ┌─────────────────────────┐\n"
    "                          │       cli.ts             │\n"
    "                          │   命令行入口 & 参数解析    │\n"
    "                          └───────────┬─────────────┘\n"
    "                                      │\n"
    "                          ┌───────────▼─────────────┐\n"
    "                          │      explorer.ts         │\n"
    "                          │     核心探索引擎          │\n"
    "                          └───┬───┬───┬───┬───┬─────┘\n"
    "                              │   │   │   │   │\n"
    "              ┌───────────────┘   │   │   │   └──────────────┐\n"
    "              │                   │   │   │                  │\n"
    "     ┌────────▼──────┐  ┌────────▼───▼───▼─────────┐  ┌─────▼──────────┐\n"
    "     │state-machine  │  │    llm-strategist.ts     │  │  backtracker   │\n"
    "     │   .ts         │  │   增强 LLM 策略层         │  │     .ts        │\n"
    "     │  页面状态机    │  └───────────┬─────────────┘  │  状态回溯管理   │\n"
    "     └───────┬───────┘              │                └───────┬─────────┘\n"
    "             │             ┌────────▼───────┐               │\n"
    "     ┌───────▼───────┐     │   callLlm()    │      ┌────────▼────────┐\n"
    "     │ form-analyzer │     │   (已有,增强)    │      │ anomaly-detector│\n"
    "     │    .ts        │     └────────────────┘      │     .ts         │\n"
    "     │  表单结构解析  │                             │  异常检测器      │\n"
    "     └───────────────┘                             └────────┬────────┘\n"
    "                                                            │\n"
    "                                                    ┌───────▼────────┐\n"
    "                                                    │ dom-snapshot   │\n"
    "                                                    │     .ts        │\n"
    "                                                    │  DOM 快照捕获  │\n"
    "                                                    └────────────────┘"
)
add_code_block(slide, Inches(1.0), Inches(1.4), Inches(11), Inches(5.8), arch, font_size=12)

# ═══════════════════════════════════════════
# Slide 11: 数据流
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_LIGHT)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), COLOR_PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "整体架构  单步执行数据流", font_size=32, color=COLOR_WHITE, bold=True)

dataflow = (
    "  Step N 开始\n"
    "     │\n"
    "     ├──→ capturePageContext()   ──→  截图 (已有)\n"
    "     │\n"
    "     ├──→ DomSnapshot.capture()  ──→  DOM 快照 JSON  [新增]\n"
    "     │\n"
    "     ├──→ StateMachine.fingerprint()  ──→  页面指纹  [新增]\n"
    "     │    │\n"
    "     │    └──→ StateMachine.classify()  ──→  PageState  [新增]\n"
    "     │         │\n"
    "     │         └──→ StateMachine.transition()  记录状态转移  [新增]\n"
    "     │\n"
    "     ├──→ FormAnalyzer.analyze()  ──→  FormSchema[]  [新增]\n"
    "     │\n"
    "     ├──→ LlmStrategist.buildSmartPrompt()  ──→  增强 Prompt  [新增]\n"
    "     │    │\n"
    "     │    └──→ callLlm()  ──→  LLM 动作  (已有, Prompt增强)\n"
    "     │\n"
    "     ├──→ executeAction()  ──→  执行动作  (已有)\n"
    "     │\n"
    "     └──→ AnomalyDetector.check()  ──→  异常?  [新增]\n"
    "          │\n"
    "          ├── YES ──→ Backtracker.recover()  [新增]\n"
    "          │            │\n"
    "          │            ├── goBack() + 重试\n"
    "          │            └── 超限 → 跳过继续\n"
    "          │\n"
    "          └── NO  ──→  Step N+1"
)
add_code_block(slide, Inches(1.5), Inches(1.3), Inches(10), Inches(6.0), dataflow, font_size=12)

# ═══════════════════════════════════════════
# Slide 12: 实施计划
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_LIGHT)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), COLOR_PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "06  实施计划与产出", font_size=32, color=COLOR_WHITE, bold=True)

steps = [
    ("Step 1", "更新 types.ts", "新增 StateMachine、FormSchema、DomSnapshot、AnomalyRecord、BacktrackSnapshot 等类型定义"),
    ("Step 2", "创建 state-machine.ts", "页面指纹计算、FSM 状态分类与转移记录"),
    ("Step 3", "创建 form-analyzer.ts", "表单结构解析、字段分组、约束检测"),
    ("Step 4", "创建 dom-snapshot.ts", "DOM 简化树捕获、JSON 序列化、前后步 diff"),
    ("Step 5", "创建 anomaly-detector.ts", "JS 错误监听、页面崩溃检测、HTTP 错误识别"),
    ("Step 6", "创建 backtracker.ts", "状态快照管理、goBack 恢复、回溯计数限制"),
    ("Step 7", "创建 llm-strategist.ts", "增强 Prompt 构建、多步规划解析"),
    ("Step 8", "更新 explorer.ts", "集成所有新模块到主循环"),
    ("Step 9", "更新 cli.ts", "新增 --backtrack、--snapshot 等 CLI 参数"),
    ("Step 10", "测试验证", "npm run build 编译通过 + 实际网页探索测试"),
]
for i, (num, title, desc) in enumerate(steps):
    row = i // 2
    col = i % 2
    left = Inches(0.6 + col * 6.2)
    top = Inches(1.5 + row * 1.1)
    color = COLOR_SECONDARY if i < 3 else COLOR_ACCENT if i < 6 else RGBColor(0x8B, 0x5C, 0xF6)
    add_step_card(slide, left, top, Inches(5.9), Inches(0.95), num, title, desc, color)

# ═══════════════════════════════════════════
# Slide 13: 预期成果 & 总结
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), COLOR_SECONDARY)
add_rect(slide, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), COLOR_SECONDARY)

add_text_box(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8),
             "📊 预期成果对比", font_size=36, color=COLOR_WHITE, bold=True)

# Before vs After
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(4.5), Inches(0.5),
             "完善前", font_size=24, color=COLOR_RED, bold=True)
before_items = [
    "随机探索，无方向感",
    "LLM 不理解表单结构",
    "无异常感知能力",
    "出错后直接跳过",
    "无法回溯到良好状态",
    "DOM 信息仅存内存",
]
add_bullet_list(slide, Inches(1.5), Inches(2.5), Inches(4.5), Inches(3.5),
                before_items, font_size=15, color=RGBColor(0x94, 0xA3, 0xB8))

add_text_box(slide, Inches(7.5), Inches(1.8), Inches(4.5), Inches(0.5),
             "完善后", font_size=24, color=COLOR_SECONDARY, bold=True)
after_items = [
    "FSM 引导，有目标探索",
    "LLM 精准理解表单语义",
    "全方位异常实时检测",
    "异常触发自动回溯恢复",
    "DOM 快照支持离线分析",
    "完整状态转移图可视化",
]
add_bullet_list(slide, Inches(7.5), Inches(2.5), Inches(4.5), Inches(3.5),
                after_items, font_size=15, color=COLOR_WHITE)

add_rect(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.03), COLOR_SECONDARY)

add_text_box(slide, Inches(1.5), Inches(5.9), Inches(10), Inches(0.8),
             "🎯 核心目标：将随机探索器升级为具备状态感知、结构理解、异常恢复能力的智能 GUI 测试工具",
             font_size=20, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 产出文件
add_text_box(slide, Inches(3.0), Inches(6.5), Inches(7), Inches(0.5),
             "新增 5 个模块文件 | 修改 3 个现有文件 | 约 800+ 行 TypeScript 代码",
             font_size=16, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# Slide 14: 修改前后详细对比
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_LIGHT)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), COLOR_PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "07  修改前后详细对比", font_size=32, color=COLOR_WHITE, bold=True)

# 上半部分：文件结构对比
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(0.4),
             "📁 修改前（3 个源文件）", font_size=18, color=COLOR_RED, bold=True)
before_code = (
    "src/\n"
    "├── cli.ts          (74 行)\n"
    "├── explorer.ts     (380 行)\n"
    "└── types.ts        (67 行)\n"
    "\n"
    "总计: 3 文件, ~520 行"
)
add_code_block(slide, Inches(0.8), Inches(1.75), Inches(5.8), Inches(2.2), before_code, font_size=11)

add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.4),
             "📁 修改后（9 个源文件）", font_size=18, color=COLOR_SECONDARY, bold=True)
after_code = (
    "src/\n"
    "├── cli.ts              [修改] 80 行\n"
    "├── explorer.ts         [修改] 550 行\n"
    "├── types.ts            [修改] 195 行\n"
    "├── state-machine.ts    [新增] 180 行\n"
    "├── form-analyzer.ts    [新增] 230 行\n"
    "├── llm-strategist.ts   [新增] 220 行\n"
    "├── dom-snapshot.ts     [新增] 200 行\n"
    "├── anomaly-detector.ts [新增] 190 行\n"
    "└── backtracker.ts      [新增] 160 行\n"
    "\n"
    "总计: 9 文件, ~2000 行"
)
add_code_block(slide, Inches(7.0), Inches(1.75), Inches(5.8), Inches(2.2), after_code, font_size=11)

# 下半部分：能力对比表
add_text_box(slide, Inches(0.8), Inches(4.1), Inches(11.5), Inches(0.4),
             "📊 能力维度对比", font_size=18, color=COLOR_PRIMARY, bold=True)

# 表头
table_y = Inches(4.55)
col_widths = [Inches(3.5), Inches(4.2), Inches(4.2)]
col_starts = [Inches(0.8), Inches(4.3), Inches(8.5)]
headers = ["能力维度", "修改前", "修改后"]
for j, (hdr, start) in enumerate(zip(headers, col_starts)):
    cell = add_rect(slide, start, table_y, col_widths[j], Inches(0.4), COLOR_PRIMARY)
    tf = cell.text_frame
    tf.paragraphs[0].text = hdr
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

comparison_rows = [
    ("页面状态感知", "❌ 无", "✅ FSM 8 种状态分类"),
    ("表单结构理解", "❌ 仅枚举 input", "✅ label 关联/required/类型推断"),
    ("LLM Prompt 质量", "基础元素列表", "状态+表单结构+填充进度+异常"),
    ("DOM 持久化", "❌ 仅内存", "✅ 每步 JSON 快照 + diff"),
    ("异常检测", "❌ 无", "✅ 6 类异常 (JS/console/crash/HTTP/空白/导航)"),
    ("故障恢复", "❌ 出错即停止", "✅ 自动回溯 goBack + 重试"),
    ("CLI 可控性", "5 个参数", "9 个参数 (含异常/回溯/快照开关)"),
]

for i, (dim, before, after) in enumerate(comparison_rows):
    row_y = table_y + Inches(0.4) + Inches(i * 0.38)
    bg_color = COLOR_WHITE if i % 2 == 0 else RGBColor(0xF1, 0xF5, 0xF9)
    for j, (text, start, width) in enumerate(zip([dim, before, after], col_starts, col_widths)):
        cell = add_rect(slide, start, row_y, width, Inches(0.38), bg_color)
        tf = cell.text_frame
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = COLOR_DARK if j == 0 else (COLOR_SECONDARY if "✅" in text else COLOR_RED)
        tf.paragraphs[0].font.name = "Microsoft YaHei"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

# ═══════════════════════════════════════════
# Slide 15: 实际完成情况
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_LIGHT)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), COLOR_SECONDARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "08  实际完成情况", font_size=32, color=COLOR_WHITE, bold=True)

# 三步完成状态
status_cards = [
    ("方向一", "页面状态抽象与表单感知", "✅ 已完成", COLOR_SECONDARY, [
        "state-machine.ts: 指纹计算 + 8 状态分类 + 转移记录",
        "form-analyzer.ts: <form> 解析 + label 关联 + required 检测",
        "types.ts: 新增 PageFingerprint/StateTransition/FormSchema 等 12 个类型",
        "explorer.ts 集成: 每步计算指纹、分类状态、记录转移",
    ]),
    ("方向二", "LLM 智能策略增强", "✅ 已完成", COLOR_ACCENT, [
        "llm-strategist.ts: buildSystemPrompt + buildUserPrompt",
        "Prompt 增强: 表单结构 → 表格输出、填充进度、状态感知",
        "异常感知: 将异常摘要注入 prompt 辅助 LLM 决策",
        "support: 兼容 OpenAI / DeepSeek / Moonshot 等 API",
    ]),
    ("方向三", "截图/DOM快照/异常/回溯", "✅ 已完成", RGBColor(0x8B, 0x5C, 0xF6), [
        "dom-snapshot.ts: 简化 DOM 树 → JSON + step diff",
        "anomaly-detector.ts: 6 类异常监听 + 严重程度判断",
        "backtracker.ts: 快照栈 + goBack/导航 + 次数限制",
        "explorer.ts 集成: 每步快照 → 检测 → 严重异常触发回溯",
    ]),
]

for i, (name, title, status, color, items) in enumerate(status_cards):
    left = Inches(0.4 + i * 4.2)
    # 卡片头
    header = add_rect(slide, left, Inches(1.4), Inches(4.0), Inches(0.55), color)
    tf = header.text_frame
    tf.paragraphs[0].text = f"{name}: {title}"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 状态标签
    add_text_box(slide, left + Inches(0.1), Inches(2.0), Inches(3.8), Inches(0.3),
                 status, font_size=14, color=color, bold=True)
    # 详情
    add_bullet_list(slide, left + Inches(0.1), Inches(2.35), Inches(3.8), Inches(3.5),
                    items, font_size=11, spacing=Pt(5))

# 底部统计
add_rect(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.04), COLOR_PRIMARY)
stats = [
    "✅ TypeScript 编译: 零错误通过 (tsc --noEmit)",
    "✅ 9 个 .js 文件成功生成到 dist/",
    "✅ 新增 CLI 参数: --anomaly-detection / --max-backtrack / --dom-snapshot / --snapshot-dir",
    "✅ 向下兼容: 所有原有参数和默认行为保持不变",
]
add_bullet_list(slide, Inches(0.8), Inches(5.95), Inches(11.7), Inches(1.2),
                stats, font_size=13, color=COLOR_DARK, spacing=Pt(3))

# ═══════════════════════════════════════════
# Slide 16: 后续完善与展望
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_LIGHT)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), COLOR_ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
             "09  后续完善与展望", font_size=32, color=COLOR_WHITE, bold=True)

# 三大展望
future_cards = [
    ("🔬", "短期（1-2 周）", COLOR_SECONDARY, [
        "状态转移图可视化: 将 FSM transitions 输出为 Mermaid/Graphviz",
        "HTML 报告生成: 将截图+快照+日志整合为可浏览的测试报告",
        "表单填充策略库: 预置常见网站（登录/注册/搜索）的填充模板",
        "A/B 对比: 同一网站多次运行，对比状态覆盖率",
    ]),
    ("🚀", "中期（1-2 月）", COLOR_ACCENT, [
        "覆盖率指标: 页面状态覆盖率、表单字段覆盖率、路径覆盖率",
        "并发探索: 多个浏览器实例并行探索不同 URL 路径",
        "断言系统: 用户可定义预期状态转移路径，自动验证",
        "Replay 回放: 基于 DOM 快照离线回放整个探索过程",
    ]),
    ("🌟", "长期（3-6 月）", RGBColor(0x8B, 0x5C, 0xF6), [
        "强化学习探索: 用 RL 替代 LLM+随机策略，自动学习最优探索路径",
        "视觉回归测试: 截图对比检测 UI 异常变更",
        "CI/CD 集成: GitHub Actions / Jenkins 插件",
        "自然语言测试: \"帮我测试登录功能\" → 自动生成探索计划",
    ]),
]

for i, (icon, title, color, items) in enumerate(future_cards):
    left = Inches(0.4 + i * 4.2)
    add_rect(slide, left, Inches(1.4), Inches(4.0), Inches(5.1), COLOR_WHITE)
    # 标题条
    title_bar = add_rect(slide, left, Inches(1.4), Inches(4.0), Inches(0.55), color)
    tf = title_bar.text_frame
    tf.paragraphs[0].text = f"{icon} {title}"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_bullet_list(slide, left + Inches(0.1), Inches(2.1), Inches(3.8), Inches(4.2),
                    items, font_size=11, spacing=Pt(6))

# ═══════════════════════════════════════════
# Slide 17: 技术亮点 & 总结
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), COLOR_ACCENT)
add_rect(slide, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), COLOR_ACCENT)

add_text_box(slide, Inches(1.5), Inches(0.6), Inches(10), Inches(0.7),
             "✨ 技术亮点总结", font_size=36, color=COLOR_WHITE, bold=True)

highlights = [
    ("🎯 页面指纹算法", "URL 路径 + 表单签名 hash + 关键元素文本 hash，精确识别页面状态"),
    ("🧩 模块化架构", "6 个独立模块，职责清晰，可单独测试，可插拔启用/禁用"),
    ("🔄 异常闭环", "检测 → 记录 → 回溯 → 恢复 → 继续，形成完整容错链路"),
    ("🧠 LLM 上下文增强", "Markdown 表格输出表单结构、进度追踪、异常感知，最大化 LLM 理解力"),
    ("📸 全链路可追溯", "截图 + DOM JSON + 日志 + 状态转移，事后可完整复盘"),
    ("🛡️ 向下兼容", "所有新功能默认启用但可关闭，原有 API 和参数完全保留"),
]

for i, (icon, desc) in enumerate(highlights):
    y = Inches(1.55 + i * 0.85)
    # 图标
    add_text_box(slide, Inches(1.5), y, Inches(0.5), Inches(0.5),
                 icon, font_size=24, color=COLOR_ACCENT, bold=True)
    # 描述
    add_text_box(slide, Inches(2.1), y + Inches(0.08), Inches(9.5), Inches(0.6),
                 desc, font_size=15, color=COLOR_WHITE)

add_rect(slide, Inches(1.5), Inches(6.15), Inches(10), Inches(0.03), COLOR_ACCENT)

add_text_box(slide, Inches(1.5), Inches(6.35), Inches(10), Inches(0.8),
             "从「随机盲探」到「智能感知」—— 完成了 Web GUI 自动测试工具的质变升级",
             font_size=22, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ── 保存 ──
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(os.path.dirname(output_dir), "完善方案讲解.pptx")
prs.save(output_path)
print(f"✅ PPT 已生成: {output_path}")
